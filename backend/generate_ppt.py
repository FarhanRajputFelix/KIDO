from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    
    # Custom Theme Colors
    PRIMARY = RGBColor(108, 99, 255)  # Purple
    SECONDARY = RGBColor(255, 107, 107) # Coral
    TEXT_COLOR = RGBColor(45, 45, 63)
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "KIDO Learning Ecosystem"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    title.text_frame.paragraphs[0].font.bold = True
    subtitle.text = "Safe, Intelligent & Interactive Digital World for Children\nBy Farhan (FelixX-Tech)"
    
    # 2. Problem Statement Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "The Problem"
    title_shape.text_frame.paragraphs[0].font.color.rgb = SECONDARY
    tf = body_shape.text_frame
    tf.text = "Current platforms fail children and parents:"
    p = tf.add_paragraph()
    p.text = "Unsafe content algorithms targeting kids"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Lack of genuine educational value"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Zero parent visibility into learning progress"
    p.level = 1
    
    # 3. The Solution Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Our Solution: KIDO"
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    tf = body_shape.text_frame
    tf.text = "An ecosystem built entirely around safety & growth."
    p = tf.add_paragraph()
    p.text = "Gamified Learning (XP, Levels, Badges)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Verified Videos without Ads or Comments"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Parent-Approved Social Network"
    p.level = 1
    
    # 4. System Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = "System Architecture"
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    tf = body_shape.text_frame
    tf.text = "Full Stack Modern Technology:"
    p = tf.add_paragraph()
    p.text = "Frontend: Flutter 3.x (Glassmorphism, High FPS)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Backend: Python FastAPI (Asynchronous, High Concurrency)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Database: PostgreSQL (16 Relational Tables)"
    p.level = 1
    
    # 5. AI Integation
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = "AI Integration"
    title_shape.text_frame.paragraphs[0].font.color.rgb = SECONDARY
    tf = body_shape.text_frame
    tf.text = "Smart engines driving KIDO:"
    p = tf.add_paragraph()
    p.text = "Adaptive Difficulty Engine (Matches UI to child's skill)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Behavioral Stress Analysis (Alerts parents on drops in engagement)"
    p.level = 1
    
    # Output path
    output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "docs", "presentation", "KIDO_Product_Presentation.pptx")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Presentation generated at: {output_path}")

if __name__ == "__main__":
    create_presentation()
